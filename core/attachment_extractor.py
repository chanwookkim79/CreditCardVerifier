"""첨부파일(PDF/PPT/PPTX/XLS/XLSX/DOC/DOCX)에서 영수증 이미지를 추출한다.

지원 형식:
  - PDF  : 각 페이지를 고해상도 PNG로 렌더링 (pypdfium2, 이미 설치됨)
  - PPTX : 슬라이드 내 embedded 이미지 추출 (python-pptx)
  - XLSX : 시트 내 embedded 이미지 추출 (openpyxl)
  - DOCX : 문서 내 embedded 이미지 추출 (python-docx)
  - DOC/PPT/XLS (구형): LibreOffice가 설치된 경우 DOCX/PPTX/XLSX 변환 후 처리,
                        없으면 NotSupportedError 발생

반환값: 추출된 이미지 임시 파일 경로 리스트 (list[Path])
       호출 측에서 사용 후 cleanup_extracted() 또는 직접 삭제할 것.
"""
import io
import tempfile
from pathlib import Path

from core.logger import get_logger

log = get_logger("attachment_extractor")

IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}
SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".xlsx", ".docx", ".ppt", ".xls", ".doc"}


class NotSupportedError(Exception):
    pass


def extract_images(file_path: str | Path) -> list[Path]:
    """첨부파일에서 이미지를 추출해 임시 파일로 저장하고 경로 리스트를 반환."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_SUFFIXES:
        raise NotSupportedError(f"지원하지 않는 파일 형식: {suffix}")

    log.info(f"[EXTRACT] {path.name} ({suffix}) 이미지 추출 시작")

    # 구형 바이너리 포맷은 LibreOffice로 변환 후 재귀 처리
    if suffix in (".doc", ".ppt", ".xls"):
        converted = _convert_via_libreoffice(path)
        results = extract_images(converted)
        converted.unlink(missing_ok=True)
        return results

    if suffix == ".pdf":
        images = _extract_from_pdf(path)
    elif suffix == ".pptx":
        images = _extract_from_pptx(path)
    elif suffix == ".xlsx":
        images = _extract_from_xlsx(path)
    elif suffix == ".docx":
        images = _extract_from_docx(path)
    else:
        images = []

    log.info(f"[EXTRACT] {path.name} → {len(images)}개 이미지 추출 완료")
    return images


def cleanup_extracted(image_paths: list[Path]) -> None:
    """extract_images()가 생성한 임시 파일 삭제."""
    for p in image_paths:
        p.unlink(missing_ok=True)


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_from_pdf(path: Path) -> list[Path]:
    import pypdfium2 as pdfium

    tmp_dir = Path(tempfile.mkdtemp(prefix="ccv_pdf_"))
    results = []

    pdf = pdfium.PdfDocument(str(path))
    for i, page in enumerate(pdf):
        # scale=2 → 144 DPI (기본 72 DPI의 2배), OCR 정확도 향상
        bitmap = page.render(scale=2)
        pil_image = bitmap.to_pil()
        out_path = tmp_dir / f"page_{i+1:03d}.png"
        pil_image.save(str(out_path), format="PNG")
        results.append(out_path)
        log.debug(f"[EXTRACT][PDF] 페이지 {i+1} → {out_path.name}")

    pdf.close()
    return results


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _extract_from_pptx(path: Path) -> list[Path]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    tmp_dir = Path(tempfile.mkdtemp(prefix="ccv_pptx_"))
    results = []
    idx = 1

    prs = Presentation(str(path))
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_blob = shape.image.blob
                ext = _blob_to_ext(shape.image.content_type)
                out_path = tmp_dir / f"slide{slide_no:03d}_img{idx:03d}{ext}"
                out_path.write_bytes(image_blob)
                results.append(out_path)
                log.debug(f"[EXTRACT][PPTX] 슬라이드 {slide_no} 이미지 {idx} → {out_path.name}")
                idx += 1

    return results


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _extract_from_xlsx(path: Path) -> list[Path]:
    from openpyxl import load_workbook
    from PIL import Image

    tmp_dir = Path(tempfile.mkdtemp(prefix="ccv_xlsx_"))
    results = []
    idx = 1

    wb = load_workbook(str(path), data_only=True)
    for sheet in wb.worksheets:
        for img_obj in sheet._images:
            try:
                raw = img_obj.ref
                data = raw.getvalue() if hasattr(raw, "getvalue") else raw
                pil_img = Image.open(io.BytesIO(data))
                out_path = tmp_dir / f"sheet_{sheet.title[:10]}_img{idx:03d}.png"
                pil_img.save(str(out_path), format="PNG")
                results.append(out_path)
                log.debug(f"[EXTRACT][XLSX] 시트 '{sheet.title}' 이미지 {idx} → {out_path.name}")
                idx += 1
            except Exception as e:
                log.warning(f"[EXTRACT][XLSX] 이미지 {idx} 추출 실패: {e}")

    return results


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _extract_from_docx(path: Path) -> list[Path]:
    from docx import Document
    from PIL import Image

    tmp_dir = Path(tempfile.mkdtemp(prefix="ccv_docx_"))
    results = []
    idx = 1

    doc = Document(str(path))
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                image_part = rel.target_part
                data = image_part.blob
                ext = Path(image_part.partname).suffix.lower()
                if ext not in IMAGE_SUFFIXES:
                    ext = ".png"
                pil_img = Image.open(io.BytesIO(data))
                out_path = tmp_dir / f"img{idx:03d}{ext}"
                pil_img.save(str(out_path), format="PNG")
                results.append(out_path)
                log.debug(f"[EXTRACT][DOCX] 이미지 {idx} → {out_path.name}")
                idx += 1
            except Exception as e:
                log.warning(f"[EXTRACT][DOCX] 이미지 {idx} 추출 실패: {e}")

    return results


# ── 구형 포맷 변환 (LibreOffice) ─────────────────────────────────────────────

def _convert_via_libreoffice(path: Path) -> Path:
    """LibreOffice headless로 doc/ppt/xls → docx/pptx/xlsx 변환."""
    import subprocess

    ext_map = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}
    target_ext = ext_map[path.suffix.lower()]
    tmp_dir = Path(tempfile.mkdtemp(prefix="ccv_lo_"))

    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to",
             target_ext.lstrip("."), "--outdir", str(tmp_dir), str(path)],
            check=True, capture_output=True, timeout=120,
        )
    except FileNotFoundError:
        raise NotSupportedError(
            f"구형 포맷({path.suffix}) 변환을 위해 LibreOffice 설치가 필요합니다. "
            f"설치: https://www.libreoffice.org/download/"
        )
    except subprocess.CalledProcessError as e:
        raise NotSupportedError(f"LibreOffice 변환 실패: {e.stderr.decode()}")

    converted = tmp_dir / (path.stem + target_ext)
    if not converted.exists():
        raise NotSupportedError(f"LibreOffice 변환 결과 파일을 찾을 수 없음: {converted}")

    log.info(f"[EXTRACT] LibreOffice 변환 완료: {path.name} → {converted.name}")
    return converted


# ── 유틸 ──────────────────────────────────────────────────────────────────────

def _blob_to_ext(content_type: str) -> str:
    """MIME type → 파일 확장자 (기본 .png)"""
    mapping = {
        "image/jpeg": ".jpg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
    }
    return mapping.get(content_type, ".png")
